from pptx import Presentation
from pptx.util import Pt
from PIL import Image
import os
import tkinter as tk
from tkinter import filedialog

def get_image_metadata(image_path):
    image = Image.open(image_path)
    metadata = image.info
    return str(metadata)

def create_presentation(image_folder, template_path):
    # Load the template
    presentation = Presentation(template_path)

    for root, dirs, files in os.walk(image_folder):
        for filename in files:
            if filename.endswith(".png"):
                slide_layout = presentation.slide_layouts[3]
                slide = presentation.slides.add_slide(slide_layout)

                title_placeholder = slide.placeholders[0]
                title_placeholder.text = os.path.splitext(filename)[0]

                content_placeholder_1 = slide.placeholders[1]
                img_path = os.path.join(root, filename)
                content_placeholder_1.element.getparent().remove(content_placeholder_1.element)
                slide.shapes.add_picture(img_path, content_placeholder_1.left, content_placeholder_1.top, content_placeholder_1.width, content_placeholder_1.height)

                metadata = get_image_metadata(os.path.join(root, filename))

                content_placeholder_2 = slide.placeholders[2]
                p = content_placeholder_2.text_frame.add_paragraph()
                p.text = "Metadata: " + metadata
                p.font.size = Pt(8)

    presentation.save(os.path.join(image_folder, "presentation.pptx"))

def browse_directory():
    global image_folder
    image_folder = filedialog.askdirectory() 

def browse_template():
    global template_path
    template_path = filedialog.askopenfilename()

def run_script():
    create_presentation(image_folder, template_path)

root = tk.Tk()
button_browse_directory = tk.Button(root, text="Select Directory", command=browse_directory)
button_browse_directory.pack()
button_browse_template = tk.Button(root, text="Select Blank PPT File", command=browse_template)
button_browse_template.pack()
button_run = tk.Button(root, text="Run", command=run_script)
button_run.pack()
root.mainloop()
